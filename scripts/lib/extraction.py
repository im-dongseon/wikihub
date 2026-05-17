"""binary → text 추출 dispatch (F2 wiki-schema.md §A3 + F3 §2.6 정본).

각 형식의 dispatch:
- binary (.pptx/.docx/.xlsx/.pdf): Python 라이브러리로 텍스트 추출
- Google native (.gdoc/.gsheet/.gslides): gws drive files export — body 는 sync.py 가 호출
- 텍스트 (.md/.txt): passthrough

추출 실패 시 wiki 페이지는 작성 (file_map 정합성) — body 만 ``[extraction failed: <reason>]``.
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Callable


@dataclass
class ExtractionResult:
    body_text: str
    tool: str
    tool_version: str
    extraction_status: str  # 'success' | 'failed'
    reason: str = ""


def _safe_version(distribution_name: str) -> str:
    """L2·SIG-4 fix — importlib.metadata 로 PyPI 패키지 버전 조회.

    Args:
        distribution_name: PyPI 배포 패키지명 (예: 'python-pptx', 'pdfminer.six').
            import path 가 아님 (예: 'pptx' 가 아니라 'python-pptx').
    """
    try:
        from importlib.metadata import version  # type: ignore
        return str(version(distribution_name))
    except Exception:
        return "unknown"


def _failed(tool: str, reason: str) -> ExtractionResult:
    return ExtractionResult(
        body_text=f"[extraction failed: {reason}]",
        tool=tool,
        tool_version=_safe_version(tool) if tool else "n/a",
        extraction_status="failed",
        reason=reason,
    )


# ---------------------------------------------------------------------------
# binary extractors
# ---------------------------------------------------------------------------

def extract_pptx(path: Path) -> ExtractionResult:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:
        return _failed("python-pptx", f"import: {e}")
    try:
        prs = Presentation(str(path))
        lines: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            lines.append(f"## Slide {i}")
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        lines.append(text)
            lines.append("")
        return ExtractionResult(
            body_text="\n".join(lines).rstrip() + "\n",
            tool="python-pptx",
            tool_version=_safe_version("python-pptx"),
            extraction_status="success",
        )
    except Exception as e:
        return _failed("python-pptx", f"runtime: {type(e).__name__}: {e}")


def extract_docx(path: Path) -> ExtractionResult:
    try:
        from docx import Document  # type: ignore
    except ImportError as e:
        return _failed("python-docx", f"import: {e}")
    try:
        doc = Document(str(path))
        lines: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = (para.style.name if para.style else "") or ""
            if style.startswith("Heading") or style.startswith("제목"):
                lines.append(f"## {text}")
            else:
                lines.append(text)
        return ExtractionResult(
            body_text="\n".join(lines) + "\n" if lines else "",
            tool="python-docx",
            tool_version=_safe_version("python-docx"),
            extraction_status="success",
        )
    except Exception as e:
        return _failed("python-docx", f"runtime: {type(e).__name__}: {e}")


def extract_xlsx(path: Path) -> ExtractionResult:
    try:
        import openpyxl  # type: ignore
    except ImportError as e:
        return _failed("openpyxl", f"import: {e}")
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        lines: list[str] = []
        for ws in wb.worksheets:
            lines.append(f"## Sheet: {ws.title}")
            headers: list[str] = []
            header_row_idx: int | None = None
            for i, row in enumerate(ws.iter_rows(max_row=5, values_only=True), start=1):
                non_empty = [v for v in row if v is not None and str(v).strip()]
                if non_empty:
                    headers = [str(v) if v is not None else "" for v in row]
                    header_row_idx = i
                    break
            if not headers:
                lines.append("(empty sheet or no header detected)")
                lines.append("")
                continue
            row_count = (ws.max_row or header_row_idx) - header_row_idx
            lines.append("headers: " + ", ".join(headers))
            lines.append(f"row count (estimate): {row_count}")
            lines.append("")
        wb.close()
        return ExtractionResult(
            body_text="\n".join(lines).rstrip() + "\n",
            tool="openpyxl",
            tool_version=_safe_version("openpyxl"),
            extraction_status="success",
        )
    except Exception as e:
        return _failed("openpyxl", f"runtime: {type(e).__name__}: {e}")


def extract_pdf(path: Path) -> ExtractionResult:
    try:
        from pdfminer.high_level import extract_text as _pdfminer_extract_text  # type: ignore
        import pdfminer  # type: ignore
    except ImportError as e:
        return _failed("pdfminer.six", f"import: {e}")
    try:
        text = _pdfminer_extract_text(str(path))
        return ExtractionResult(
            body_text=text.strip() + "\n" if text else "",
            tool="pdfminer.six",
            tool_version=_safe_version("pdfminer.six"),
            extraction_status="success",
        )
    except Exception as e:
        return _failed("pdfminer.six", f"runtime: {type(e).__name__}: {e}")


def extract_text(path: Path) -> ExtractionResult:
    """passthrough — .md / .txt 그대로."""
    try:
        return ExtractionResult(
            body_text=path.read_text(encoding="utf-8"),
            tool="passthrough",
            tool_version="n/a",
            extraction_status="success",
        )
    except UnicodeDecodeError:
        try:
            return ExtractionResult(
                body_text=path.read_text(encoding="cp949"),
                tool="passthrough",
                tool_version="n/a",
                extraction_status="success",
            )
        except Exception as e:
            return _failed("passthrough", f"decode: {e}")
    except Exception as e:
        return _failed("passthrough", f"read: {type(e).__name__}: {e}")


# ---------------------------------------------------------------------------
# MIME dispatch — F2 wiki-schema.md §A3 정합
# ---------------------------------------------------------------------------

# V<N> Phase 2 결함 #9 fix (2026-05-17 — ADR-0027 Q1 lock):
# rclone mount 가 `--drive-export-formats docx,xlsx,pptx,md` 우선순위로 Google native 를
# binary export → mount path 는 `.docx`·`.xlsx`·`.pptx` 확장자 포함. sync.py 의 mount
# lookup 시점에서만 mimeType→ext suffix 적용 후 본 dispatch 로 binary 변환.
LOCAL_EXTRACTION_DISPATCH: dict[str, Callable[[Path], ExtractionResult]] = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_pptx,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_docx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": extract_xlsx,
    "application/pdf": extract_pdf,
    "text/markdown": extract_text,
    "text/plain": extract_text,
    "text/csv": extract_text,
    "application/vnd.google-apps.document": extract_docx,
    "application/vnd.google-apps.spreadsheet": extract_xlsx,
    "application/vnd.google-apps.presentation": extract_pptx,
}


# Google native MIME → rclone mount 가 export 하는 binary MIME 매핑 (ADR-0027 Q1 lock).
GWS_EXPORT_MIME: dict[str, str] = {
    "application/vnd.google-apps.document": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.google-apps.spreadsheet": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.google-apps.presentation": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def extract(path: Path, mime_type: str) -> ExtractionResult:
    """Local 파일 + MIME → ExtractionResult.

    Google native 는 본 함수가 처리 안 함 (sync.py 가 gws export 후 ``extract_text`` 호출).
    매핑 없는 MIME 은 failed 반환.
    """
    handler = LOCAL_EXTRACTION_DISPATCH.get(mime_type)
    if handler is None:
        return _failed("dispatch", f"unsupported MIME: {mime_type}")
    return handler(path)


__all__ = [
    "ExtractionResult",
    "LOCAL_EXTRACTION_DISPATCH",
    "GWS_EXPORT_MIME",
    "extract",
    "extract_pptx",
    "extract_docx",
    "extract_xlsx",
    "extract_pdf",
    "extract_text",
]
